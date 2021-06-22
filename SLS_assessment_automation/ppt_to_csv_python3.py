# NOTE: To execute the program, you need to install the python-pptx module. To do so, go to your console and type "python -m pip install python-pptx"
import csv
from pptx import Presentation
# IMPORTANT: Enter the location of the file here
prs = Presentation("data/D2.2_Decimal items.pptx");

# a dictionary indicating the cell position of the information in the tables
LABELS = ["SN", "level_1", "level_2", "KU_label", "KU_desc", "item_no", "pre_req", "question", "answer", "option_1", "option_2", "option_3", "option_4"]
CELL_LOCATION = {
	"SN": [0, 0],
	"level_1": [1, 1],
	"level_2": [1, 2],
	"KU_label": [3, 1],
	"KU_desc": [4, 1],
	"item_no": [5, 1],
	"pre_req": [6, 1],
	"question": [1, 3],
	"answer": [1, 4]
}

# interating through each individual slide
data_arr = []
for i, slide in enumerate(prs.slides):
	data = {}
	for shape in slide.shapes:
		# if the current shape has a table 	
		if shape.has_table:
			cur_table = shape.table
			# extracting the information from the tables
			for key in CELL_LOCATION:
				row = CELL_LOCATION[key][0]
				col = CELL_LOCATION[key][1]

				data[key] = cur_table.cell(row, col).text.strip()

			# attempting to extract the options from the question
			# NOTE: This method requires that there are no breaklines within the actual question itself. (If not, part of the question will appear in the options)
			options = data["question"].split("\n")
			opt_cnt = 1

			for opt_idx, opt in enumerate(options):
				# do not use the first element, as it is the actual question itself
				if opt_idx == 0:
					continue
				opt = opt.strip()
				# if opt is not blank, add it to the data
				if opt != '':
					opt_str = "option_{}".format(opt_cnt)
					data[opt_str] = opt
					opt_cnt += 1

			# if this question is multiple choice, remove the options from the question string
			if opt_cnt > 1:
				data["question"] = options[0]


	# displaying for clarity
	print("Slide {}\n--------------------".format(i + 1))
	for key in data:
		print("{}: {}".format(key, data[key]))
	print("\n")

	data_arr.append(data)

# writing to external csv file
with open("output.csv", mode="w", newline='', encoding='utf-8') as output:
	writer = csv.writer(output)
	# adding the titles
	writer.writerow(LABELS)
	for data in data_arr:	
		# if no data recorded
		if len(data.values()) == 0:
			continue
		# the purpose behind this is to ensure that each row will have the same number of columns
		row = []
		for label in LABELS:
			if label in data:
				row.append(data[label])
			else:
				row.append('')
		writer.writerow(row)



